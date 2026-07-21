"""Generates the pitch deck. Run: python deck/generate_deck.py
Requires docs/architecture.png and deck/screenshots/*.png to be present first.
"""
from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.font.size = Pt(20)
    return slide


def add_image_slide(prs, title, image_path, width=Inches(9)):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=width)
    return slide


def main():
    prs = Presentation()

    add_title_slide(
        prs, "Throughline",
        "Cross-Channel Journey Stitching — the visibility layer every other AmEx agent theme's data flows through",
    )

    add_bullet_slide(prs, "The Problem", [
        "A dispute starts on the app, escalates through a call, resolves in-branch",
        "Today: four disconnected systems, zero unified view of that one journey",
        "No way to see where it broke down, or that it broke down at all",
    ])

    add_bullet_slide(prs, "Why This Theme", [
        "Not a 7th agent competing with the other 6 CodeStreet themes",
        "It's the visibility layer that makes any of those journeys traceable end to end",
        "Every other theme generates cross-channel data; none of it is stitched today",
    ])

    add_image_slide(prs, "Architecture", "docs/architecture.png")

    add_image_slide(prs, "Live Demo — Journey Timeline", "deck/screenshots/journey_timeline.png")
    add_image_slide(prs, "Live Demo — Aggregate Patterns", "deck/screenshots/aggregate_patterns.png")

    add_bullet_slide(prs, "Measured Results", [
        "100% identity resolution accuracy on the seeded set (incl. ambiguous cases correctly left unresolved)",
        "6.28ms average pipeline latency per event, ingestion to resolved timeline",
        "High-friction customers return 5x less often, measured against our seeded synthetic dataset",
    ])

    add_bullet_slide(prs, "Business Impact", [
        "Turns four blind systems into one traceable customer story",
        "Surfaces churn risk from journey friction before the customer leaves",
        "Honest by design: low-confidence links are flagged, never force-matched",
    ])

    add_bullet_slide(prs, "Scalability & What's Next", [
        "Event generators -> real Kafka/Spark streaming ingestion at scale",
        "SQLite -> Snowflake/BigQuery for the event store",
        "Analytics layer output feeds a downstream tool like Amplitude/Mixpanel",
    ])

    prs.save("deck/Throughline_Deck.pptx")
    print("Saved deck/Throughline_Deck.pptx")


if __name__ == "__main__":
    main()
